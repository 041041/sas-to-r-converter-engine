import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank layout

    # Colors
    NAVY_DARK = RGBColor(11, 19, 43)      # #0B132B Title & Hero Dark
    SLATE_DARK = RGBColor(26, 29, 36)    # #1A1D24 Code editor / Dark Cards
    CARD_DARK = RGBColor(20, 27, 45)     # #141B2D Dark card background
    LIGHT_BG = RGBColor(248, 250, 252)   # #F8FAFC Light Slide Background
    WHITE = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(15, 23, 42)     # #0F172A Title/Main text
    TEXT_MUTED = RGBColor(71, 85, 105)   # #475569 Subtitles/Details
    TEXT_LIGHT = RGBColor(241, 245, 249) # #F1F5F9 Light text on dark
    TEXT_LIGHT_MUTED = RGBColor(148, 163, 184)
    
    ACCENT_TEAL = RGBColor(14, 165, 233)  # #0EA5E9 Electric Teal
    ACCENT_GREEN = RGBColor(16, 185, 129) # #10B981 Mint Emerald
    ACCENT_BLUE = RGBColor(59, 130, 246)  # #3B82F6 Royal Blue
    ACCENT_PURPLE = RGBColor(124, 58, 237)# #7C3AED Vivid Purple
    BORDER_COLOR = RGBColor(226, 232, 240)
    BORDER_DARK = RGBColor(51, 65, 85)

    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, subtitle_text, dark=False):
        # Header text frame
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.1))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = WHITE if dark else TEXT_DARK
        
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(13)
        p2.font.name = "Arial"
        p2.font.color.rgb = TEXT_LIGHT_MUTED if dark else TEXT_MUTED
        p2.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 1 — TITLE SLIDE (Dark Navy Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1, NAVY_DARK)

    # Decorative top accent bar
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_TEAL
    top_bar.line.fill.background()

    # Title box
    tx1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf1 = tx1.text_frame
    tf1.word_wrap = True

    # Category Pill Tag
    p_tag = tf1.paragraphs[0]
    p_tag.text = "CLINICAL DATA SCIENCE & SAS MODERNIZATION PLATFORM"
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.name = "Arial"
    p_tag.font.color.rgb = ACCENT_TEAL
    p_tag.space_after = Pt(16)

    # Main Title
    p_main = tf1.add_paragraph()
    p_main.text = "SAS to R Modernization Studio"
    p_main.font.size = Pt(40)
    p_main.font.bold = True
    p_main.font.name = "Arial"
    p_main.font.color.rgb = WHITE
    p_main.space_after = Pt(12)

    # Subtitle
    p_sub = tf1.add_paragraph()
    p_sub.text = "AI-Assisted Clinical SAS Migration & Validation"
    p_sub.font.size = Pt(20)
    p_sub.font.name = "Arial"
    p_sub.font.color.rgb = TEXT_LIGHT_MUTED

    # Bottom Metadata Cards (3 horizontal pills)
    features = [
        ("LLM-Assisted Translation", "Intelligent code analysis & R generation"),
        ("Macro Resolution", "Automated dependency expansion & parsing"),
        ("Closed-Loop R Validation", "Rscript execution & dataset comparison")
    ]
    card_w = Inches(3.6)
    card_h = Inches(1.1)
    top_pos = Inches(5.4)

    for i, (feat_title, feat_desc) in enumerate(features):
        left_pos = Inches(1.0 + i * 3.86)
        card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = BORDER_DARK
        
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_right = Inches(0.2)
        tf_c.margin_top = Inches(0.18)
        
        p = tf_c.paragraphs[0]
        p.text = feat_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL
        
        p_d = tf_c.add_paragraph()
        p_d.text = feat_desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_LIGHT_MUTED
        p_d.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 2 — SAS TO R MODERNIZATION CHALLENGE
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2, LIGHT_BG)
    add_header(slide2, "SAS to R Modernization Challenge", "Key Operational & Technical Friction Points in Clinical Migration")

    challenges = [
        ("Legacy SAS Codebase Overhead", "Decades of accumulated SAS DATA steps, PROCs, and custom macros require modernization without breaking validated analytical pipelines.", ACCENT_BLUE),
        ("Macro Dependency Chains", "Deeply nested macros, global %let variables, and dynamic code generation obscure data lineage and manual conversion logic.", ACCENT_PURPLE),
        ("Manual Rewrite Effort", "Line-by-line manual translation from SAS to R consumes thousands of bio-statistical hours, introducing high costs and slow timelines.", ACCENT_TEAL),
        ("Quality & Traceability Risks", "Manual translation carries substantial risk of subtle calculation errors, requiring rigorous side-by-side validation against baseline outputs.", ACCENT_GREEN)
    ]

    for i, (ch_title, ch_desc, border_c) in enumerate(challenges):
        col = i % 2
        row = i // 2
        left_pos = Inches(0.8 + col * 5.95)
        top_pos = Inches(1.8 + row * 2.5)
        
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(5.75), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)

        # Color accent bar on left of card
        bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, Inches(0.12), Inches(2.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = border_c
        bar.line.fill.background()

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)
        
        p = tf.paragraphs[0]
        p.text = ch_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        p_d = tf.add_paragraph()
        p_d.text = ch_desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 3 — OUR SOLUTION: END-TO-END MODERNIZATION PLATFORM
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3, LIGHT_BG)
    add_header(slide3, "Our Solution: SAS-to-R Modernization Platform", "Automated Workflow from Macro Resolution to Execution Verification")

    # Workflow Steps
    steps = [
        ("1. SAS Program", "Import legacy SAS scripts & dependencies"),
        ("2. Macro Resolution", "Expand parameters & resolve nested macros"),
        ("3. SAS Parsing", "Extract DATA steps, PROCs & logic tree"),
        ("4. SAS → R Conversion", "LLM-assisted conversion to Modern R / Base R"),
        ("5. R Validation", "Execute Rscript & compare DataFrame output"),
        ("6. Conversion Audit", "Diagnostic trace & execution status report")
    ]

    box_w = Inches(1.75)
    box_h = Inches(2.2)
    gap = Inches(0.2)

    for i, (st_title, st_desc) in enumerate(steps):
        left_pos = Inches(0.8 + i * (1.75 + 0.2))
        top_pos = Inches(2.0)

        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = ACCENT_TEAL if i == 3 else BORDER_COLOR
        box.line.width = Pt(2) if i == 3 else Pt(1)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = st_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL if i == 3 else TEXT_DARK
        p.alignment = PP_ALIGN.CENTER

        p_d = tf.add_paragraph()
        p_d.text = st_desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(8)
        p_d.alignment = PP_ALIGN.CENTER

        # Draw arrow between boxes
        if i < len(steps) - 1:
            arr_left = left_pos + box_w + Inches(0.04)
            arr = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arr_left, top_pos + Inches(0.9), Inches(0.12), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT_BLUE
            arr.line.fill.background()

    # Solution Highlights Footer Card
    foot_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(11.733), Inches(2.1))
    foot_card.fill.solid()
    foot_card.fill.fore_color.rgb = CARD_DARK
    foot_card.line.color.rgb = BORDER_DARK

    tf_f = foot_card.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_right = Inches(0.3)
    tf_f.margin_top = Inches(0.25)

    p = tf_f.paragraphs[0]
    p.text = "CORE PLATFORM PRINCIPLES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL

    points = [
        "• Preserves Analytical Intent: Accurately maps SAS DATA step expressions, PROCs, and implicit loops into structured R functions.",
        "• Configurable R Dialects: Generates clean, readable Modern R (tidyverse) or pure Base R based on organizational standards.",
        "• Closed-Loop Verification: Automatically executes generated R code against sample data to verify shape, types, and values."
    ]
    for pt in points:
        p_pt = tf_f.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_LIGHT
        p_pt.space_before = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 4 — SAS → R CONVERSION (MAIN SLIDE)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4, NAVY_DARK)
    add_header(slide4, "SAS → R Conversion", "Representative Demonstration: DATA Step & PROC SQL Aggregation", dark=True)

    # Top Mini Workflow Bar
    wf_items = ["SAS Program", "Macro Resolution", "SAS Understanding", "R Translation", "Execution Check", "Final R Code"]
    wf_w = Inches(1.8)
    wf_gap = Inches(0.18)
    for i, item in enumerate(wf_items):
        l_pos = Inches(0.8 + i * (wf_w + wf_gap))
        t_pos = Inches(1.6)
        
        pill = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_pos, t_pos, wf_w, Inches(0.45))
        pill.fill.solid()
        pill.fill.fore_color.rgb = ACCENT_TEAL if item == "R Translation" else CARD_DARK
        pill.line.color.rgb = ACCENT_TEAL if item == "R Translation" else BORDER_DARK
        
        tf_p = pill.text_frame
        tf_p.word_wrap = True
        tf_p.margin_top = Inches(0.08)
        p = tf_p.paragraphs[0]
        p.text = item
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE if item == "R Translation" else TEXT_LIGHT_MUTED
        p.alignment = PP_ALIGN.CENTER

    # Code Snippets Side-by-Side
    # Left: SAS Code Box
    sas_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.25), Inches(5.7), Inches(4.7))
    sas_box.fill.solid()
    sas_box.fill.fore_color.rgb = SLATE_DARK
    sas_box.line.color.rgb = BORDER_DARK

    tf_sas = sas_box.text_frame
    tf_sas.word_wrap = True
    tf_sas.margin_left = tf_sas.margin_right = Inches(0.25)
    tf_sas.margin_top = Inches(0.2)

    p = tf_sas.paragraphs[0]
    p.text = "ORIGINAL SAS CODE (orders.sas)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(8)

    sas_code = (
        "/* 1. DATA Step & Datalines Creation */\n"
        "data orders;\n"
        "  input cust_id $ order_id amount;\n"
        "  datalines;\n"
        "  C101 1001 250\n"
        "  C101 1002 400\n"
        "  C102 1003 150\n"
        "  C101 1004 300\n"
        "  C103 1005 500\n"
        ";\n"
        "run;\n\n"
        "/* 2. PROC SQL Aggregation & Filtering */\n"
        "proc sql;\n"
        "  create table customer_summary as\n"
        "  select cust_id,\n"
        "         count(order_id) as total_orders,\n"
        "         sum(amount) as total_spent\n"
        "  from orders\n"
        "  group by cust_id\n"
        "  having total_spent > 500\n"
        "  order by total_spent desc;\n"
        "quit;"
    )
    p_code = tf_sas.add_paragraph()
    p_code.text = sas_code
    p_code.font.size = Pt(10.5)
    p_code.font.name = "Consolas"
    p_code.font.color.rgb = TEXT_LIGHT

    # Right: Generated R Code Box
    r_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.25), Inches(5.733), Inches(4.7))
    r_box.fill.solid()
    r_box.fill.fore_color.rgb = SLATE_DARK
    r_box.line.color.rgb = ACCENT_GREEN
    r_box.line.width = Pt(1.5)

    tf_r = r_box.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = Inches(0.25)
    tf_r.margin_top = Inches(0.2)

    p = tf_r.paragraphs[0]
    p.text = "GENERATED MODERN R (tidyverse)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(8)

    r_code = (
        "# 1. Data Frame Initialization\n"
        "orders <- data.frame(\n"
        '  cust_id = c("C101", "C101", "C102", "C101", "C103"),\n'
        "  order_id = c(1001, 1002, 1003, 1004, 1005),\n"
        "  amount = c(250, 400, 150, 300, 500)\n"
        ")\n\n"
        "# 2. PROC SQL Translation Pipeline\n"
        "customer_summary <- orders %>%\n"
        "  group_by(cust_id) %>%\n"
        "  summarise(\n"
        "    total_orders = n(),\n"
        "    total_spent = sum(amount),\n"
        '    .groups = "drop"\n'
        "  ) %>%\n"
        "  filter(total_spent > 500) %>%\n"
        "  arrange(desc(total_spent))\n\n"
        "customer_summary"
    )
    p_rcode = tf_r.add_paragraph()
    p_rcode.text = r_code
    p_rcode.font.size = Pt(10.5)
    p_rcode.font.name = "Consolas"
    p_rcode.font.color.rgb = TEXT_LIGHT

    # -------------------------------------------------------------
    # SLIDE 5 — CONVERSION CAPABILITIES
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5, LIGHT_BG)
    add_header(slide5, "Intelligent Conversion Capabilities", "Verified SAS Constructs & Automated R Translation Rules")

    caps = [
        ("DATA Step Translation", "SET, IF/ELSE statements, conditional variable creation, and column assignments."),
        ("PROC SQL Conversion", "SELECT, JOIN, WHERE filters, GROUP BY aggregations, HAVING clauses, and ORDER BY."),
        ("PROC SORT & Logic", "BY-variable sorting, DESCENDING keywords, and FIRST./LAST. group indexing."),
        ("PROC MEANS & FREQ", "Summary statistics, cross-tabulations (count & rename), and frequency distributions."),
        ("PROC REPORT & TRANSPOSE", "Wide/long dataset reshaping, summary subtotals, and column structure mapping."),
        ("SAS Function Hints", "Built-in mapping for INTCK, INTNX, COMPRESS, CATX, SCAN, SUBSTR, PUT, and INPUT.")
    ]

    for i, (cap_title, cap_desc) in enumerate(caps):
        col = i % 3
        row = i // 3
        l_pos = Inches(0.8 + col * 3.97)
        t_pos = Inches(1.8 + row * 2.2)

        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_pos, t_pos, Inches(3.78), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_COLOR

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = cap_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        p_d = tf.add_paragraph()
        p_d.text = cap_desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(6)

    # Capability Disclaimer Card (Bottom)
    disc_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.9))
    disc_card.fill.solid()
    disc_card.fill.fore_color.rgb = CARD_DARK
    disc_card.line.color.rgb = BORDER_DARK

    tf_disc = disc_card.text_frame
    tf_disc.word_wrap = True
    tf_disc.margin_left = Inches(0.3)
    tf_disc.margin_top = Inches(0.18)

    p_disc = tf_disc.paragraphs[0]
    p_disc.text = "CAPABILITY SCOPE & TRANSPARENCY"
    p_disc.font.size = Pt(10)
    p_disc.font.bold = True
    p_disc.font.color.rgb = ACCENT_TEAL

    p_dt = tf_disc.add_paragraph()
    p_dt.text = "“Automates supported SAS constructs and identifies complex cases requiring review.”"
    p_dt.font.size = Pt(12)
    p_dt.font.bold = True
    p_dt.font.color.rgb = WHITE
    p_dt.space_before = Pt(2)

    # -------------------------------------------------------------
    # SLIDE 6 — MACRO & COMPLEX SAS HANDLING
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6, LIGHT_BG)
    add_header(slide6, "Macro & Complex SAS Handling", "Automated Parsing & Expansion of SAS Macro Dependencies")

    # Left Box: Visual Dependency Chain
    chain_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.1))
    chain_card.fill.solid()
    chain_card.fill.fore_color.rgb = CARD_DARK
    chain_card.line.color.rgb = BORDER_DARK

    tf_ch = chain_card.text_frame
    tf_ch.word_wrap = True
    tf_ch.margin_left = Inches(0.3)
    tf_ch.margin_top = Inches(0.3)

    p = tf_ch.paragraphs[0]
    p.text = "MACRO RESOLUTION TREE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(14)

    nodes = [
        ("MAIN_PROGRAM.SAS", "Master execution script"),
        ("↓  %include / Macro Call", ""),
        ("%MACRO_A(study=CDISC)", "Global parameter extraction"),
        ("↓  Nested Invocation", ""),
        ("%MACRO_B(ds=adsl)", "Table & variable scoping"),
        ("↓  Child Utility Call", ""),
        ("%MACRO_C(var=age)", "Statement expansion & code inline")
    ]

    for node_title, node_desc in nodes:
        p_n = tf_ch.add_paragraph()
        if "↓" in node_title:
            p_n.text = node_title
            p_n.font.size = Pt(11)
            p_n.font.bold = True
            p_n.font.color.rgb = ACCENT_BLUE
            p_n.space_before = Pt(2)
        else:
            p_n.text = f"• {node_title}"
            p_n.font.size = Pt(12)
            p_n.font.bold = True
            p_n.font.color.rgb = WHITE
            p_n.space_before = Pt(4)
            if node_desc:
                p_sub = tf_ch.add_paragraph()
                p_sub.text = f"   {node_desc}"
                p_sub.font.size = Pt(10)
                p_sub.font.color.rgb = TEXT_LIGHT_MUTED

    # Right Box: Macro Handling Capabilities
    macro_caps = [
        ("Macro Discovery & Parsing", "Automatically scans SAS libraries to extract %macro definitions, parameter lists, and default values."),
        ("Global %LET Substitution", "Tracks global %let variables across multiple files and substitutes values into downstream DATA steps."),
        ("Nested Macro Resolution", "Recursively resolves multi-level macro calls, expanding macro bodies into plain executable SAS code."),
        ("Context-Aware Parameter Mapping", "Maps &parameter references to R function arguments or data frame column selectors."),
        ("Resolution Diagnostics", "Generates warnings for unresolved macro variables, missing includes, or circular macro dependencies.")
    ]

    for i, (mc_title, mc_desc) in enumerate(macro_caps):
        t_pos = Inches(1.8 + i * 1.02)
        m_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), t_pos, Inches(6.933), Inches(0.92))
        m_card.fill.solid()
        m_card.fill.fore_color.rgb = WHITE
        m_card.line.color.rgb = BORDER_COLOR

        tf_m = m_card.text_frame
        tf_m.word_wrap = True
        tf_m.margin_left = Inches(0.25)
        tf_m.margin_top = Inches(0.12)

        p = tf_m.paragraphs[0]
        p.text = mc_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        p_d = tf_m.add_paragraph()
        p_d.text = mc_desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(2)

    # -------------------------------------------------------------
    # SLIDE 7 — QUALITY & VALIDATION
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7, LIGHT_BG)
    add_header(slide7, "Quality & Execution Validation", "Closed-Loop R Script Verification & Output Comparison")

    # Workflow bar
    v_items = ["1. SAS Input", "2. Generated R", "3. R Execution", "4. Output Validation", "5. Quality Report"]
    v_w = Inches(2.15)
    for i, v_item in enumerate(v_items):
        l_pos = Inches(0.8 + i * (2.15 + 0.24))
        card_v = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_pos, Inches(1.8), v_w, Inches(0.7))
        card_v.fill.solid()
        card_v.fill.fore_color.rgb = CARD_DARK
        card_v.line.color.rgb = BORDER_DARK

        tf_v = card_v.text_frame
        tf_v.word_wrap = True
        tf_v.margin_top = Inches(0.2)
        p = tf_v.paragraphs[0]
        p.text = v_item
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL
        p.alignment = PP_ALIGN.CENTER

    # Middle Section: Actual Implemented Validation Indicators (User Request)
    val_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.7), Inches(11.733), Inches(2.2))
    val_card.fill.solid()
    val_card.fill.fore_color.rgb = WHITE
    val_card.line.color.rgb = ACCENT_TEAL
    val_card.line.width = Pt(1.5)

    tf_vc = val_card.text_frame
    tf_vc.word_wrap = True
    tf_vc.margin_left = Inches(0.3)
    tf_vc.margin_top = Inches(0.2)

    p = tf_vc.paragraphs[0]
    p.text = "EXECUTION & VALIDATION METRICS ENGINE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

    # 4 Status Badges Grid
    status_boxes = [
        ("Conversion Status", "Completed", ACCENT_GREEN),
        ("Execution Status", "Validated", ACCENT_BLUE),
        ("Conversion Warnings", "Displayed when applicable", ACCENT_TEAL),
        ("Manual Review", "Flagged when required", ACCENT_PURPLE)
    ]

    for i, (sb_label, sb_val, sb_col) in enumerate(status_boxes):
        l_pos = Inches(1.1 + i * 2.85)
        s_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_pos, Inches(3.4), Inches(2.6), Inches(1.2))
        s_box.fill.solid()
        s_box.fill.fore_color.rgb = CARD_DARK
        s_box.line.color.rgb = sb_col

        tf_s = s_box.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = Inches(0.15)
        tf_s.margin_top = Inches(0.18)

        p_lbl = tf_s.paragraphs[0]
        p_lbl.text = sb_label
        p_lbl.font.size = Pt(10)
        p_lbl.font.color.rgb = TEXT_LIGHT_MUTED
        p_lbl.alignment = PP_ALIGN.CENTER

        p_v = tf_s.add_paragraph()
        p_v.text = sb_val
        p_v.font.size = Pt(12)
        p_v.font.bold = True
        p_v.font.color.rgb = sb_col
        p_v.space_before = Pt(4)
        p_v.alignment = PP_ALIGN.CENTER

    # Bottom Details Card
    det_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.1), Inches(11.733), Inches(1.8))
    det_card.fill.solid()
    det_card.fill.fore_color.rgb = CARD_DARK
    det_card.line.color.rgb = BORDER_DARK

    tf_det = det_card.text_frame
    tf_det.word_wrap = True
    tf_det.margin_left = Inches(0.3)
    tf_det.margin_top = Inches(0.2)

    p = tf_det.paragraphs[0]
    p.text = "KEY VALIDATION CAPABILITIES"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL

    val_points = [
        "• Rscript Subprocess Execution: Runs generated R script in an isolated environment to verify syntax and library imports.",
        "• DataFrame Comparison Engine: Evaluates row dimensions, column names, character case, and numeric values within a set tolerance (1e-3).",
        "• Automated Feedback Loop: Feeds execution runtime errors back to the translation engine for automated R script repair."
    ]
    for vp in val_points:
        p_vp = tf_det.add_paragraph()
        p_vp.text = vp
        p_vp.font.size = Pt(11)
        p_vp.font.color.rgb = TEXT_LIGHT
        p_vp.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 8 — SUPPORTING CAPABILITIES
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8, LIGHT_BG)
    add_header(slide8, "Supporting Platform Capabilities", "Integrated Modules Enhancing Clinical Migration Workflows")

    supp_cards = [
        ("PROJECT WORKSPACE", "Organize SAS programs, macro libraries, clinical metadata, generated R scripts, and audit logs within a structured environment.", ACCENT_BLUE),
        ("CLINICAL METADATA", "Support clinical data standard metadata workflows, variable mapping specifications, and Define.xml structures.", ACCENT_PURPLE),
        ("TLF SUPPORT", "Interpret clinical TLF (Tables, Listings, Figures) shell specifications to generate structured R code utilizing gt and text formatting.", ACCENT_TEAL),
        ("ANALYTICS / GRAPHING", "Interactive Graph Builder module to create clinical charts (Bar, Line, Scatter, Boxplot) powered by ggplot2 with custom styling.", ACCENT_GREEN)
    ]

    for i, (sc_title, sc_desc, sc_color) in enumerate(supp_cards):
        col = i % 2
        row = i // 2
        l_pos = Inches(0.8 + col * 5.95)
        t_pos = Inches(1.8 + row * 2.5)

        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_pos, t_pos, Inches(5.75), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)

        # Top Accent Strip
        strip = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, l_pos, t_pos, Inches(5.75), Inches(0.1))
        strip.fill.solid()
        strip.fill.fore_color.rgb = sc_color
        strip.line.fill.background()

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = sc_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        p_d = tf.add_paragraph()
        p_d.text = sc_desc
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 9 — CLIENT VALUE & NEXT STEPS
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9, NAVY_DARK)
    add_header(slide9, "Client Value & Next Steps", "Strategic Outcomes for Clinical SAS Modernization", dark=True)

    # Value Flow Cards (Top)
    outcomes = [
        ("Accelerate Modernization", "Dramatically shorten SAS-to-R migration timelines."),
        ("Reduce Manual Effort", "Minimize repetitive manual rewriting & bio-statistical costs."),
        ("Improve Consistency", "Standardize R output to modern tidyverse / Base R patterns."),
        ("Early Risk Detection", "Identify conversion gaps & logic edge cases automatically."),
        ("Full Traceability", "Maintain transparent conversion trace & side-by-side audit logs.")
    ]

    o_w = Inches(2.15)
    o_gap = Inches(0.24)
    for i, (o_title, o_desc) in enumerate(outcomes):
        l_pos = Inches(0.8 + i * (o_w + o_gap))
        o_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_pos, Inches(1.8), o_w, Inches(2.4))
        o_card.fill.solid()
        o_card.fill.fore_color.rgb = CARD_DARK
        o_card.line.color.rgb = BORDER_DARK

        tf_o = o_card.text_frame
        tf_o.word_wrap = True
        tf_o.margin_left = tf_o.margin_right = Inches(0.15)
        tf_o.margin_top = Inches(0.2)

        p = tf_o.paragraphs[0]
        p.text = f"0{i+1}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL
        p.alignment = PP_ALIGN.CENTER

        p_t = tf_o.add_paragraph()
        p_t.text = o_title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_before = Pt(6)
        p_t.alignment = PP_ALIGN.CENTER

        p_d = tf_o.add_paragraph()
        p_d.text = o_desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_LIGHT_MUTED
        p_d.space_before = Pt(6)
        p_d.alignment = PP_ALIGN.CENTER

    # Next Steps Hero Banner (Bottom)
    ns_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.3))
    ns_card.fill.solid()
    ns_card.fill.fore_color.rgb = CARD_DARK
    ns_card.line.color.rgb = ACCENT_GREEN
    ns_card.line.width = Pt(2)

    tf_ns = ns_card.text_frame
    tf_ns.word_wrap = True
    tf_ns.margin_left = Inches(0.4)
    tf_ns.margin_top = Inches(0.25)

    p = tf_ns.paragraphs[0]
    p.text = "RECOMMENDED NEXT STEP FOR CLIENT ENGAGEMENT"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(8)

    p_call = tf_ns.add_paragraph()
    p_call.text = "“Validate the platform against representative client SAS programs and macro libraries.”"
    p_call.font.size = Pt(18)
    p_call.font.bold = True
    p_call.font.color.rgb = WHITE
    p_call.space_after = Pt(10)

    ns_steps = [
        "1. Pilot Run: Execute modernization studio on 5–10 representative client SAS programs.",
        "2. Custom Macro Calibration: Integrate client-specific macro libraries into resolution processor.",
        "3. Validation Benchmark: Deliver complete side-by-side output comparison and execution reports."
    ]
    for nss in ns_steps:
        p_s = tf_ns.add_paragraph()
        p_s.text = nss
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = TEXT_LIGHT_MUTED
        p_s.space_before = Pt(2)

    # Save presentation
    output_path = "/Users/sandeep/.gemini/antigravity/scratch/sas-to-r-converter/SAS_to_R_Modernization_Studio.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
